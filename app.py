#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Web版个性化学习多智能体系统
FastAPI + SSE + 完整多智能体工作流 + 精美UI
"""

import os
import io as _io
import json
import re as _re
import asyncio
import requests
import subprocess
import tempfile
from urllib.parse import quote
from fastapi import FastAPI, Request, HTTPException, UploadFile, File, Form
from fastapi.responses import HTMLResponse, FileResponse, Response
from fastapi.staticfiles import StaticFiles
from sse_starlette.sse import EventSourceResponse
from dotenv import load_dotenv

from docx import Document as _Document
from docx.shared import Inches as _Inches, Pt as _Pt, RGBColor as _RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH as _WD_ALIGN
from docx.oxml.ns import qn as _qn

try:
    from pypdf import PdfReader as _PdfReader
except ImportError:
    _PdfReader = None

from graph_web import process_message
from video_downloader import download_video_async, get_download_progress, get_video_info_simple
from database import list_sessions, get_session, get_session_messages, delete_session, \
    insert_document, list_documents, get_document, delete_document, upsert_session, get_plan, \
    get_profile, save_profile, save_plan as db_save_plan, get_documents_content
from core import build_ppt_prompt, parse_ppt_json, call_llm_sync, \
    build_video_prompt, parse_video_script, sanitize_video_scenes, \
    render_video_frames, generate_video_narration, compose_video, \
    build_video_fallback_script, call_doubao_image, \
    build_plan_prompt, build_profile_prompt, extract_json, merge_profile

load_dotenv()

app = FastAPI(title="AI智学 - 个性化学习多智能体系统")

# 挂载静态文件目录
static_dir = os.path.join(os.path.dirname(__file__), "static")
if os.path.isdir(static_dir):
    app.mount("/static", StaticFiles(directory=static_dir), name="static")

# 从 index.html 加载前端页面
_INDEX_HTML = None


def _load_index_html() -> str:
    html_path = os.path.join(os.path.dirname(__file__), "index.html")
    if os.path.exists(html_path):
        with open(html_path, "r", encoding="utf-8") as f:
            return f.read()
    return "<html><body><h1>前端页面丢失</h1></body></html>"


@app.get("/", response_class=HTMLResponse)
async def index():
    return HTMLResponse(
        content=_load_index_html(),
        headers={
            "Cache-Control": "no-cache, no-store, must-revalidate",
            "Pragma": "no-cache",
            "Expires": "0",
        },
    )


def _decode_base64_file(content_b64: str, ext: str) -> str:
    """将前端发来的 base64 文件内容解码并解析为文本。"""
    import base64
    try:
        raw = base64.b64decode(content_b64)
    except Exception:
        raise HTTPException(status_code=400, detail="文件 base64 解码失败")
    if ext == ".pdf":
        return _parse_pdf(raw)
    elif ext == ".docx":
        return _parse_docx(raw)
    elif ext == ".doc":
        return _parse_doc(raw)
    else:
        # fallback: try UTF-8
        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return raw.decode("gbk")
            except UnicodeDecodeError:
                raise HTTPException(status_code=400, detail="无法解码文件内容")


def _build_file_context(files: list) -> str:
    """将前端 files 数组格式化为消息前缀。"""
    lines = ["[用户上传了以下文件]\n"]
    for i, f in enumerate(files):
        fname = f.get("name", f"file_{i}")
        ftype = f.get("type", "text")
        ext = f.get("ext", "")
        content = f.get("content", "")
        if ftype == "image":
            lines.append(f"📷 图片 {fname}: data:{f.get('mime','image/png')};base64,{content[:80]}...")
        elif ftype == "text":
            lines.append(f"📄 文件 {fname}:\n```\n{content}\n```")
        elif ftype == "binary":
            parsed = _decode_base64_file(content, ext)
            lines.append(f"📄 文件 {fname}:\n```\n{parsed}\n```")
    lines.append("[文件内容结束]\n")
    return "\n".join(lines)


@app.post("/chat/stream")
async def chat_stream(request: Request):
    data = await request.json()
    user_message = data.get("message", "")
    session_id = data.get("session_id", "default")
    image_mode = data.get("image_mode", False)
    web_search = data.get("web_search", False)
    deep_thinking = data.get("deep_thinking", False)
    resource_type = data.get("resource_type", "")
    files = data.get("files", [])
    print(f"[DEBUG chat_stream] web_search={web_search} image_mode={image_mode} deep_thinking={deep_thinking} resource_type={resource_type} files={len(files)} msg={user_message[:40]}", flush=True)

    # 输入验证
    if (not user_message or not user_message.strip()) and not files:
        async def empty_response():
            yield {"data": json.dumps({"type": "status", "content": "⚠️ 消息不能为空，请输入学习相关问题。"})}
            yield {"data": json.dumps({"type": "end"})}
        return EventSourceResponse(empty_response())
    if len(user_message) > 2000:
        async def too_long_response():
            yield {"data": json.dumps({"type": "status", "content": "⚠️ 消息过长（超过2000字符），请精简后重试。"})}
            yield {"data": json.dumps({"type": "end"})}
        return EventSourceResponse(too_long_response())
    if len(session_id) > 128:
        async def bad_session_response():
            yield {"data": json.dumps({"type": "status", "content": "⚠️ session_id 格式无效。"})}
            yield {"data": json.dumps({"type": "end"})}
        return EventSourceResponse(bad_session_response())

    # 构建多模态消息
    full_message = user_message
    if files:
        file_context = _build_file_context(files)
        full_message = file_context + "\n" + (user_message or "请帮我分析上传的文件内容")

    async def event_generator():
        gen = process_message(session_id, full_message, image_mode=image_mode, web_search=web_search, deep_thinking=deep_thinking, resource_type=resource_type)
        ended = False
        try:
            for event_type, content in gen:
                if await request.is_disconnected():
                    # 关闭 generator 以触发其 finally 块（持久化AI消息）
                    gen.close()
                    break
                yield {"data": json.dumps({"type": event_type, "content": content})}
                if event_type == "end":
                    ended = True
        except asyncio.CancelledError:
            gen.close()
        except Exception as e:
            gen.close()
            yield {"data": json.dumps({"type": "error", "content": str(e)})}
        finally:
            if not ended:
                yield {"data": json.dumps({"type": "end"})}

    return EventSourceResponse(event_generator())


@app.post("/video/download")
async def start_video_download(request: Request):
    data = await request.json()
    bvid = data.get("bvid", "").strip()
    if not bvid:
        raise HTTPException(status_code=400, detail="请提供视频 BV 号")
    task_id = download_video_async(bvid)
    return {"task_id": task_id, "status": "started"}


@app.get("/video/download/{task_id}")
async def check_video_download(task_id: str):
    progress = get_download_progress(task_id)
    if progress.get("status") == "not_found":
        raise HTTPException(status_code=404, detail="任务不存在或已过期")
    return progress


@app.get("/video/info/{bvid}")
async def video_info(bvid: str):
    try:
        info = get_video_info_simple(bvid)
        return info
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/sessions")
async def api_list_sessions(limit: int = 20):
    return list_sessions(limit=limit)


@app.get("/sessions/{session_id}/messages")
async def api_get_session_messages(session_id: str):
    session = get_session(session_id)
    if session is None:
        raise HTTPException(status_code=404, detail="会话不存在")
    messages = get_session_messages(session_id)
    plan = get_plan(session_id)
    return {"session": session, "messages": messages, "plan": plan}


@app.delete("/sessions/{session_id}")
async def api_delete_session(session_id: str):
    session = get_session(session_id)
    if session is None:
        raise HTTPException(status_code=404, detail="会话不存在")
    delete_session(session_id)
    return {"status": "deleted", "session_id": session_id}


ALLOWED_EXTENSIONS = {
    ".txt", ".md", ".json", ".csv", ".html", ".py", ".js", ".css", ".xml", ".yaml", ".log",
    ".pdf", ".docx", ".doc",
}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB


def _parse_pdf(raw: bytes) -> str:
    """从 PDF 字节数据提取文本。"""
    if _PdfReader is None:
        raise HTTPException(status_code=500, detail="pypdf 库未安装，无法解析 PDF")
    try:
        reader = _PdfReader(_io.BytesIO(raw))
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
        return "\n\n".join(parts)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF 解析失败: {str(e)}")


def _parse_doc(raw: bytes) -> str:
    """通过 Word COM 自动化从旧版 .doc 提取文本。"""
    try:
        import pythoncom
        from win32com.client import Dispatch
    except ImportError:
        raise HTTPException(status_code=500, detail="pywin32 库未安装，无法解析 .doc 文件，请转换为 .docx 格式后重试")

    # 将 bytes 写入临时文件（COM 需要文件路径）
    tmp_path = None
    try:
        fd, tmp_path = tempfile.mkstemp(suffix=".doc")
        os.close(fd)
        with open(tmp_path, "wb") as f:
            f.write(raw)

        pythoncom.CoInitialize()
        word = Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0

        doc = word.Documents.Open(tmp_path)
        text = doc.Content.Text

        doc.Close(False)
        word.Quit()
        pythoncom.CoUninitialize()

        if not text or not text.strip():
            raise HTTPException(status_code=400, detail=".doc 文件无法提取文本，请转换为 .docx 格式后重试")
        return text
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f".doc 解析失败: {str(e)}")
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


def _parse_docx(raw: bytes) -> str:
    """从 DOCX 字节数据提取文本。"""
    if _Document is None:
        raise HTTPException(status_code=500, detail="python-docx 库未安装，无法解析 DOCX")
    try:
        doc = _Document(_io.BytesIO(raw))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        return "\n".join(parts)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DOCX 解析失败: {str(e)}")


@app.post("/knowledge/upload")
async def upload_knowledge_file(session_id: str = Form(...), file: UploadFile = File(...)):
    """上传文件导入知识库，支持 txt/md/pdf/docx 等格式。"""
    if not file.filename:
        raise HTTPException(status_code=400, detail="未选择文件")

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"不支持的文件类型: {ext}。支持: {', '.join(sorted(ALLOWED_EXTENSIONS))}")

    try:
        raw = await file.read()
    except Exception:
        raise HTTPException(status_code=400, detail="读取文件失败")

    if len(raw) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail=f"文件过大（最大 10MB），当前 {len(raw) / 1024 / 1024:.1f}MB")

    # 根据文件类型解析
    if ext == ".pdf":
        content = _parse_pdf(raw)
    elif ext == ".docx":
        content = _parse_docx(raw)
    elif ext == ".doc":
        content = _parse_doc(raw)
    else:
        try:
            content = raw.decode("utf-8")
        except UnicodeDecodeError:
            try:
                content = raw.decode("gbk")
            except UnicodeDecodeError:
                raise HTTPException(status_code=400, detail="无法解码文件内容，请使用 UTF-8 或 GBK 编码")

    # 生成预览（前 200 个字符）
    preview = content[:200].replace("\n", " ").strip()
    if len(content) > 200:
        preview += "..."

    content_type = ext.lstrip(".")
    upsert_session(session_id)  # 确保 session 存在以满足外键约束
    doc_id = insert_document(session_id, file.filename, content, content_type, len(raw))
    return {
        "status": "ok",
        "doc_id": doc_id,
        "title": file.filename,
        "size": len(raw),
        "preview": preview,
        "chars": len(content),
    }


@app.post("/knowledge/import-text")
async def import_knowledge_text(request: Request):
    """粘贴文本导入知识库。"""
    data = await request.json()
    session_id = data.get("session_id", "").strip()
    title = data.get("title", "").strip()
    content = data.get("content", "")

    if not session_id:
        raise HTTPException(status_code=400, detail="缺少 session_id")
    if not title:
        raise HTTPException(status_code=400, detail="请输入文档标题")
    if not content or not content.strip():
        raise HTTPException(status_code=400, detail="请输入文档内容")
    if len(content) > 2 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="内容过长（最大 2MB）")

    upsert_session(session_id)
    doc_id = insert_document(session_id, title, content, "text", len(content.encode("utf-8")))
    return {"status": "ok", "doc_id": doc_id, "title": title}


@app.get("/knowledge/{session_id}")
async def api_list_knowledge(session_id: str):
    """列出会话的所有知识库文档。"""
    return list_documents(session_id)


@app.delete("/knowledge/{doc_id}")
async def api_delete_knowledge(doc_id: int):
    """删除指定知识库文档。"""
    doc = get_document(doc_id)
    if doc is None:
        raise HTTPException(status_code=404, detail="文档不存在")
    delete_document(doc_id)
    return {"status": "deleted", "doc_id": doc_id}


@app.post("/resource/export-docx")
async def export_resource_docx(request: Request):
    """将 Markdown 资源内容转换为 Word 文档并下载。"""
    data = await request.json()
    markdown = data.get("content", "")
    title = data.get("title", "学习资源")

    if not markdown or not markdown.strip():
        raise HTTPException(status_code=400, detail="内容为空，无法生成文档")

    doc = _Document()

    # 设置默认字体
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Microsoft YaHei"
    font.size = _Pt(11)
    style.element.rPr.rFonts.set(_qn("w:eastAsia"), "微软雅黑")

    # 设置页边距
    for section in doc.sections:
        section.top_margin = _Inches(0.8)
        section.bottom_margin = _Inches(0.8)
        section.left_margin = _Inches(1.0)
        section.right_margin = _Inches(1.0)

    # 文档标题
    doc_title = doc.add_heading(title, level=0)
    doc_title.alignment = _WD_ALIGN.CENTER

    lines = markdown.split("\n")
    i = 0
    in_code_block = False
    code_lang = ""
    code_lines = []

    def add_formatted_paragraph(text, style_name=None):
        """添加带格式化（粗体、斜体、行内代码）的段落。"""
        p = doc.add_paragraph()
        if style_name:
            p.style = style_name
        _parse_inline(p, text)
        return p

    def _parse_inline(paragraph, text):
        """解析行内格式：**粗体**、*斜体*、`代码`。"""
        # 正则匹配行内格式
        pattern = r"(\*\*(.+?)\*\*|\*(.+?)\*|`([^`]+)`|\[([^\]]+)\]\(([^)]+)\))"
        last_end = 0
        for m in _re.finditer(pattern, text):
            # 添加前面的纯文本
            if m.start() > last_end:
                paragraph.add_run(text[last_end:m.start()])
            if m.group(2):  # **粗体**
                run = paragraph.add_run(m.group(2))
                run.bold = True
            elif m.group(3):  # *斜体*
                run = paragraph.add_run(m.group(3))
                run.italic = True
            elif m.group(4):  # `行内代码`
                run = paragraph.add_run(m.group(4))
                run.font.name = "Consolas"
                run.font.size = _Pt(10)
            elif m.group(5):  # [链接](url)
                run = paragraph.add_run(m.group(5))
                run.underline = True
                run.font.color.rgb = _RGBColor(37, 99, 235)
            last_end = m.end()
        # 添加剩余纯文本
        if last_end < len(text):
            paragraph.add_run(text[last_end:])

    def add_code_block(lang, code):
        """添加代码块。"""
        if lang:
            p = doc.add_paragraph()
            run = p.add_run(f"语言: {lang}")
            run.font.size = _Pt(9)
            run.font.color.rgb = _RGBColor(100, 100, 100)
            run.italic = True
        for line in code:
            p = doc.add_paragraph()
            p.style = doc.styles["No Spacing"] if "No Spacing" in [s.name for s in doc.styles] else doc.styles["Normal"]
            p.paragraph_format.space_before = _Pt(0)
            p.paragraph_format.space_after = _Pt(0)
            run = p.add_run(line)
            run.font.name = "Consolas"
            run.font.size = _Pt(9.5)
        doc.add_paragraph()  # 空行分隔

    def add_table_block(rows_data):
        """添加表格。"""
        if not rows_data or len(rows_data) < 2:
            return
        num_cols = max(len(row) for row in rows_data)
        table = doc.add_table(rows=len(rows_data), cols=num_cols)
        table.style = "Light Grid Accent 1"
        for r_idx, row in enumerate(rows_data):
            for c_idx, cell_text in enumerate(row):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text.strip()
                    for para in cell.paragraphs:
                        para.paragraph_format.space_before = _Pt(2)
                        para.paragraph_format.space_after = _Pt(2)
                        for run in para.runs:
                            run.font.size = _Pt(10)
                            if r_idx == 0:
                                run.bold = True
        doc.add_paragraph()  # 空行分隔

    while i < len(lines):
        line = lines[i]

        # 代码块开始
        m_code_start = _re.match(r"^```(\w*)$", line)
        if m_code_start and not in_code_block:
            in_code_block = True
            code_lang = m_code_start.group(1) or ""
            code_lines = []
            i += 1
            continue

        # 代码块结束
        if line.strip() == "```" and in_code_block:
            add_code_block(code_lang, code_lines)
            in_code_block = False
            code_lang = ""
            code_lines = []
            i += 1
            continue

        # 代码块内容
        if in_code_block:
            code_lines.append(line)
            i += 1
            continue

        # 空行
        if not line.strip():
            i += 1
            continue

        # 表格（| ... |）
        if "|" in line and line.strip().startswith("|"):
            table_rows = []
            while i < len(lines) and "|" in lines[i] and lines[i].strip().startswith("|"):
                row_line = lines[i].strip()
                if not _re.match(r"^\|[\s\-:|]+\|$", row_line):  # 跳过分隔行
                    cells = [c.strip() for c in row_line.split("|")[1:-1]]
                    if cells:
                        table_rows.append(cells)
                i += 1
            if table_rows:
                add_table_block(table_rows)
            continue

        # 标题
        heading_match = _re.match(r"^(#{1,4})\s+(.+)$", line)
        if heading_match:
            level = min(len(heading_match.group(1)), 3)  # Word 支持1-3级，映射#->1, ##->2, ###->3
            text = heading_match.group(2).strip()
            # 清除标题中的 markdown 标记
            text = _re.sub(r"\*\*(.+?)\*\*", r"\1", text)
            text = _re.sub(r"\*(.+?)\*", r"\1", text)
            text = _re.sub(r"`(.+?)`", r"\1", text)
            doc.add_heading(text, level=level)
            i += 1
            continue

        # 水平线
        if _re.match(r"^---+$", line.strip()):
            doc.add_paragraph("─" * 50)
            i += 1
            continue

        # 无序列表
        ul_match = _re.match(r"^(\s*)[-*]\s+(.+)$", line)
        if ul_match:
            text = ul_match.group(2)
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            _parse_inline(p, text)
            i += 1
            continue

        # 有序列表
        ol_match = _re.match(r"^(\s*)\d+\.\s+(.+)$", line)
        if ol_match:
            text = ol_match.group(2)
            p = doc.add_paragraph(style="List Number")
            p.clear()
            _parse_inline(p, text)
            i += 1
            continue

        # 普通段落
        add_formatted_paragraph(line.strip())
        i += 1

    # 保存到内存
    buf = _io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    safe_filename = _re.sub(r'[\\/*?:"<>|]', "_", title)[:50]
    encoded_fn = quote(f"{safe_filename}.docx")
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{encoded_fn}"}
    )


@app.post("/resource/export-pptx")
async def export_resource_pptx(request: Request):
    """将 Markdown 资源内容转换为 PPT 并下载。"""
    data = await request.json()
    markdown = data.get("content", "")
    title = data.get("title", "学习资源")

    if not markdown or not markdown.strip():
        raise HTTPException(status_code=400, detail="内容为空，无法生成PPT")

    if len(markdown) > 20000:
        raise HTTPException(status_code=400, detail="内容过长（最多20000字符）")

    # 调用 DeepSeek 生成幻灯片 JSON
    prompt = build_ppt_prompt(markdown, title)
    raw = call_llm_sync([{"role": "user", "content": prompt}])
    slides_data = parse_ppt_json(raw)

    # 降级：LLM 生成失败时，从 Markdown 手动拆分
    if not slides_data or "slides" not in slides_data:
        lines = [l.strip() for l in markdown.split("\n") if l.strip()]
        slides_data = {
            "title": title,
            "slides": [{"title": "封面", "bullets": [title], "layout": "title"}]
        }
        current = {"title": "", "bullets": [], "layout": "content"}
        for line in lines:
            if line.startswith("## "):
                if current["bullets"]:
                    slides_data["slides"].append(current)
                current = {"title": line.lstrip("# ").strip(), "bullets": [], "layout": "content"}
            elif line.startswith("# "):
                if current["bullets"]:
                    slides_data["slides"].append(current)
                current = {"title": line.lstrip("# ").strip(), "bullets": [], "layout": "content"}
            elif line.startswith("- ") or line.startswith("* "):
                current["bullets"].append(line.lstrip("-* ").strip())
            else:
                current["bullets"].append(line[:40])
        if current["bullets"]:
            slides_data["slides"].append(current)
        slides_data["slides"].insert(0, {"title": title, "bullets": ["AI智学 · 个性化学习系统"], "layout": "title"})
        slides_data["slides"].append({"title": "Q&A", "bullets": ["感谢观看", "欢迎提问"], "layout": "summary"})

    # ================================================================
    # 组装 PPT 文件 — 专业设计
    # ================================================================
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt, Emu
    from pptx.dml.color import RGBColor as PptRGB
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn as pptx_qn

    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    # ── 专业配色方案 ──
    C = {
        "primary":    PptRGB(37, 99, 235),     # 宝蓝
        "primary_dark": PptRGB(30, 64, 175),   # 深蓝
        "accent":     PptRGB(99, 102, 241),    # 靛蓝
        "accent2":    PptRGB(236, 72, 153),    # 玫红点缀
        "dark":       PptRGB(15, 23, 42),      # 近黑文字
        "body":       PptRGB(51, 65, 85),      # 正文灰
        "muted":      PptRGB(148, 163, 184),   # 浅灰
        "bg_light":   PptRGB(248, 250, 252),   # 页背景
        "card_bg":    PptRGB(255, 255, 255),   # 卡片白
        "white":      PptRGB(255, 255, 255),
        "gradient_top": PptRGB(15, 23, 42),    # 封面渐变上
        "gradient_bot": PptRGB(30, 64, 175),   # 封面渐变下
    }

    def _rect(slide, l, t, w, h, fill_color=None, border_color=None, border_width=None, radius=None):
        """添加矩形，支持圆角。"""
        shape = slide.shapes.add_shape(
            1, PptInches(l), PptInches(t), PptInches(w), PptInches(h))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = PptPt(border_width or 1)
        else:
            shape.line.fill.background()
        return shape

    def _rounded_rect(slide, l, t, w, h, fill_color, border_color=None, border_width=None):
        """添加圆角矩形（左上、右上、左下、右下统一半径）。"""
        shape = slide.shapes.add_shape(
            5, PptInches(l), PptInches(t), PptInches(w), PptInches(h))  # 5 = rounded rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = PptPt(border_width or 1)
        else:
            shape.line.fill.background()
        # 设置圆角半径
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
        return shape

    def _add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _textbox(slide, l, t, w, h, text, size=18, color=None, bold=False,
                 align=PP_ALIGN.LEFT, font="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
        """添加文本框。"""
        txBox = slide.shapes.add_textbox(PptInches(l), PptInches(t), PptInches(w), PptInches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        try:
            tf.paragraphs[0].alignment = align
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(size)
        p.font.color.rgb = color or C["body"]
        p.font.bold = bold
        p.font.name = font
        # 东亚字体
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(pptx_qn('a:altLang'), 'zh-CN')
        return tf

    def _add_rich_textbox(slide, l, t, w, h, items, color=None, size=16):
        """添加多段落文本框，每项一个段落。"""
        txBox = slide.shapes.add_textbox(PptInches(l), PptInches(t), PptInches(w), PptInches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = PptPt(size)
            p.font.color.rgb = color or C["body"]
            p.font.name = "Microsoft YaHei"
            p.space_after = PptPt(6)
        return tf

    def _circle(slide, cx, cy, r, fill_color):
        """添加实心圆。"""
        d = r * 2
        shape = slide.shapes.add_shape(
            9, PptInches(cx - r), PptInches(cy - r), PptInches(d), PptInches(d))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def _slide_number(slide, num, total):
        """右下角页码。"""
        _textbox(slide, 12.0, 7.05, 1.2, 0.35, f"{num} / {total}",
                 size=9, color=C["muted"], align=PP_ALIGN.RIGHT)

    # ── 构建幻灯片 ──
    slides_list = slides_data.get("slides", [])
    total = len(slides_list)

    for idx, slide_info in enumerate(slides_list):
        layout_type = slide_info.get("layout", "content")
        slide_title = slide_info.get("title", "")
        bullets = slide_info.get("bullets", [])
        notes_text = slide_info.get("notes", "")
        blank = prs.slide_layouts[6]

        # ═══════════════════════════════════════
        # 封面页
        # ═══════════════════════════════════════
        if layout_type == "title" or idx == 0:
            slide = prs.slides.add_slide(blank)
            _add_bg(slide, C["gradient_top"])

            # 大色块装饰
            _rect(slide, 0, 0, 13.333, 7.5, fill_color=C["gradient_top"])
            # 底部渐变条模拟
            _rect(slide, 0, 5.0, 13.333, 2.5, fill_color=C["gradient_bot"])

            # 装饰几何图形
            _circle(slide, 11.5, 1.5, 1.8, C["primary"])        # 大圆
            _circle(slide, 2.0, 5.8, 0.6, C["accent2"])          # 小玫红圆
            # 半透明方块装饰
            r1 = _rect(slide, 9.8, 0.8, 0.8, 0.8, fill_color=PptRGB(59, 130, 246))
            r1.fill.fore_color.brightness = 0.3

            # 标题
            main_title = slide_title or title
            _textbox(slide, 1.8, 2.2, 9.7, 1.6, main_title,
                     size=48, color=C["white"], bold=True, align=PP_ALIGN.LEFT)
            # 装饰线
            _rect(slide, 1.8, 3.85, 2.0, 0.06, fill_color=C["accent2"])
            # 副标题
            subtitle = bullets[0] if bullets else "个性化学习系统"
            _textbox(slide, 1.8, 4.1, 9.0, 0.8, subtitle,
                     size=22, color=PptRGB(203, 213, 225), align=PP_ALIGN.LEFT)
            # 底部信息
            _textbox(slide, 1.8, 5.8, 5.0, 0.4, "AI 智学 · 多智能体学习平台",
                     size=13, color=PptRGB(148, 163, 184), align=PP_ALIGN.LEFT)
            _textbox(slide, 1.8, 6.2, 5.0, 0.4, "Powered by DeepSeek",
                     size=10, color=C["muted"], align=PP_ALIGN.LEFT)

        # ═══════════════════════════════════════
        # 总结/Q&A 页
        # ═══════════════════════════════════════
        elif layout_type == "summary" or idx == total - 1:
            slide = prs.slides.add_slide(blank)
            _add_bg(slide, C["gradient_top"])
            _rect(slide, 0, 0, 13.333, 7.5, fill_color=C["gradient_top"])

            # 装饰
            _circle(slide, 2.0, 5.5, 1.2, C["primary"])
            _circle(slide, 11.0, 2.0, 0.5, C["accent2"])

            # 标题
            _textbox(slide, 2, 1.8, 9.3, 1.5, slide_title,
                     size=40, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
            _rect(slide, 5.5, 3.3, 2.3, 0.06, fill_color=C["accent2"])

            # 要点
            y = 3.8
            for i, b in enumerate(bullets[:5]):
                _circle(slide, 2.5, y + 0.15, 0.14, C["accent2"])
                _textbox(slide, 3.0, y, 8.3, 0.6, b,
                         size=18, color=PptRGB(203, 213, 225), align=PP_ALIGN.LEFT)
                y += 0.55

            _textbox(slide, 2, 6.5, 9.3, 0.4, "感谢观看 · AI智学出品",
                     size=12, color=C["muted"], align=PP_ALIGN.CENTER)

        # ═══════════════════════════════════════
        # 目录页
        # ═══════════════════════════════════════
        elif layout_type == "toc":
            slide = prs.slides.add_slide(blank)
            _add_bg(slide, C["white"])
            # 顶部色条
            _rect(slide, 0, 0, 13.333, 0.06, fill_color=C["primary"])
            # 标题
            _textbox(slide, 1.2, 0.5, 4.0, 0.8, slide_title or "目录",
                     size=34, color=C["dark"], bold=True)
            _rect(slide, 1.2, 1.3, 1.5, 0.04, fill_color=C["accent2"])

            y = 2.0
            for i, b in enumerate(bullets[:8]):
                # 编号圆圈
                _circle(slide, 1.8, y + 0.12, 0.28, C["primary"])
                _textbox(slide, 1.56, y - 0.02, 0.56, 0.5, str(i + 1),
                         size=14, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
                # 条目文本
                _textbox(slide, 2.5, y, 8.5, 0.5, b,
                         size=18, color=C["dark"])
                # 分隔线
                _rect(slide, 2.5, y + 0.52, 9.0, 0.01, fill_color=PptRGB(226, 232, 240))
                y += 0.65

            _slide_number(slide, idx + 1, total)

        # ═══════════════════════════════════════
        # 内容页
        # ═══════════════════════════════════════
        else:
            slide = prs.slides.add_slide(blank)
            _add_bg(slide, C["bg_light"])

            # 顶部导航条
            _rect(slide, 0, 0, 13.333, 0.06, fill_color=C["primary"])

            # 左侧色条装饰
            _rect(slide, 0, 0, 0.08, 7.5, fill_color=C["primary"])

            # 标题区域背景
            _rect(slide, 0.8, 0.35, 11.7, 1.15, fill_color=C["white"])
            _rounded_rect(slide, 0.8, 0.35, 11.7, 1.15, fill_color=C["white"])

            # 标题
            _textbox(slide, 1.3, 0.45, 10.5, 0.6, slide_title,
                     size=30, color=C["dark"], bold=True)
            _rect(slide, 1.3, 1.05, 1.8, 0.04, fill_color=C["accent2"])

            # 内容卡片
            card_l = 0.8
            card_t = 1.85
            card_w = 8.2
            card_h = min(5.0, 0.3 + len(bullets) * 0.75)

            _rounded_rect(slide, card_l, card_t, card_w, card_h,
                          fill_color=C["white"], border_color=PptRGB(226, 232, 240), border_width=0.5)

            # 右侧信息卡
            info_l = 9.4
            info_t = 1.85
            info_w = 3.3
            _rounded_rect(slide, info_l, info_t, info_w, 2.8,
                          fill_color=C["primary"], border_color=None)

            # 右侧卡内文字
            _textbox(slide, info_l + 0.3, info_t + 0.3, info_w - 0.6, 0.4, f"第 {idx + 1} 页",
                     size=11, color=PptRGB(191, 219, 254), align=PP_ALIGN.LEFT)
            _textbox(slide, info_l + 0.3, info_t + 0.7, info_w - 0.6, 0.4, f"共 {total} 页",
                     size=11, color=PptRGB(191, 219, 254), align=PP_ALIGN.LEFT)
            _textbox(slide, info_l + 0.3, info_t + 1.3, info_w - 0.6, 0.8,
                     "AI智学\n多智能体学习平台",
                     size=12, color=PptRGB(191, 219, 254), align=PP_ALIGN.LEFT)

            # 要点列表
            y = card_t + 0.35
            for i, b in enumerate(bullets[:6]):
                # 编号圆圈
                _circle(slide, card_l + 0.45, y + 0.13, 0.2, C["primary"])
                _textbox(slide, card_l + 0.22, y + 0.0, 0.46, 0.4, str(i + 1),
                         size=11, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
                # 要点文本
                _textbox(slide, card_l + 1.0, y, card_w - 1.5, 0.7, b if len(b) < 60 else b[:57] + "...",
                         size=17, color=C["body"])
                y += 0.72

            # 底部品牌条
            _rect(slide, 0, 7.1, 13.333, 0.4, fill_color=C["primary"])
            _textbox(slide, 0.8, 7.1, 5.0, 0.35, "AI智学 · 个性化学习系统",
                     size=8, color=PptRGB(191, 219, 254))
            _slide_number(slide, idx + 1, total)

        # 备注
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass

    # ── 保存 ──
    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    safe_filename = _re.sub(r'[\\/*?:"<>|]', "_", title)[:50]
    encoded_fn = quote(f"{safe_filename}.pptx")
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{encoded_fn}"}
    )


@app.post("/resource/export-mp4")
async def export_resource_mp4(request: Request):
    """根据用户问题生成教学视频 MP4。"""
    import tempfile
    import os as _os
    import shutil

    data = await request.json()
    question = data.get("question", "").strip()
    reference = data.get("reference", "")

    if not question:
        raise HTTPException(status_code=400, detail="问题为空，无法生成视频")

    # 以用户问题为核心，参考资料仅供准确性校验
    prompt = build_video_prompt(question, reference)
    raw = call_llm_sync([{"role": "user", "content": prompt}])
    scenes = parse_video_script(raw)

    # 降级：LLM 失败时手动拆分
    if not scenes:
        scenes = build_video_fallback_script(reference, question)

    # 清洗：修复封面标题 AI 幻觉
    scenes = sanitize_video_scenes(scenes, question)

    # 创建临时工作目录
    work_dir = tempfile.mkdtemp(prefix="ai_video_")

    try:
        # 渲染视频帧
        frames = render_video_frames(scenes, question, work_dir)

        # 生成旁白
        audio_files = await generate_video_narration(scenes, work_dir)

        # 合成视频
        output_path = _os.path.join(work_dir, "output.mp4")
        success = compose_video(frames, audio_files, output_path)

        if not success or not _os.path.exists(output_path):
            raise HTTPException(status_code=500, detail="视频合成失败，请检查 ffmpeg 是否安装")

        # 读取视频文件
        with open(output_path, "rb") as f:
            video_data = f.read()

        safe_filename = _re.sub(r'[\\/*?:"<>|]', "_", question)[:50]
        encoded_fn = quote(f"{safe_filename}.mp4")
        return Response(
            content=video_data,
            media_type="video/mp4",
            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{encoded_fn}"}
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"视频生成失败: {str(e)}")
    finally:
        # 清理临时文件
        try:
            shutil.rmtree(work_dir, ignore_errors=True)
        except Exception:
            pass


@app.post("/api/generate-plan")
async def generate_plan(request: Request):
    """基于历史对话深度分析，生成个性化学习路径。独立于资源生成管线。"""
    data = await request.json()
    session_id = data.get("session_id", "").strip()
    if not session_id:
        raise HTTPException(status_code=400, detail="缺少 session_id")

    # 1. 加载历史对话（优先当前会话；若为空则合并所有会话）
    messages = get_session_messages(session_id)
    if not messages:
        all_sessions = list_sessions(limit=100)
        for s in all_sessions:
            msgs = get_session_messages(s["session_id"])
            messages.extend(msgs)
        if not messages:
            raise HTTPException(status_code=400, detail="没有任何对话历史，请先进行一些学习对话")

    # 构建对话上下文（用于画像分析）
    ctx_lines = []
    for m in messages[-20:]:  # 取最近20条
        role = "用户" if m["role"] == "human" else "AI"
        content = m["content"][:300]
        ctx_lines.append(f"{role}: {content}")
    conv_ctx = "\n".join(ctx_lines)

    # 2. 加载/分析画像
    profile = get_profile(session_id)
    if not profile:
        # 首次分析画像
        prompt = build_profile_prompt({}, conv_ctx)
        resp = call_llm_sync([{"role": "user", "content": prompt}])
        new_p = extract_json(resp)
        if new_p:
            profile = merge_profile({}, new_p)
        else:
            profile = {}
    else:
        # 增量更新画像
        prompt = build_profile_prompt(profile, conv_ctx[-2000:])
        resp = call_llm_sync([{"role": "user", "content": prompt}])
        new_p = extract_json(resp)
        if new_p:
            profile = merge_profile(profile, new_p)

    # 持久化画像
    if profile:
        try:
            save_profile(session_id, profile)
        except Exception:
            pass

    # 3. 获取课程上下文
    course_ctx = ""
    try:
        docs = get_documents_content(session_id)
    except Exception:
        docs = []
    if docs:
        doc_texts = [f"[{d['title']}]: {d['content'][:500]}" for d in docs[:3]]
        course_ctx = "已上传文档:\n" + "\n".join(doc_texts)

    # 4. 生成学习路径
    plan_prompt = build_plan_prompt(profile, course_ctx)
    plan_text = call_llm_sync([{"role": "user", "content": plan_prompt}])

    parsed = extract_json(plan_text)
    if parsed and isinstance(parsed, dict) and "steps" in parsed:
        steps = parsed["steps"]
        summary = parsed.get("summary", "个性化学习路径")
        diagnosis = parsed.get("diagnosis", "")
    else:
        steps = []
        summary = "个性化学习路径"
        diagnosis = ""
        for line in plan_text.split("\n"):
            line = line.strip()
            if line and len(line) > 2 and line[0].isdigit() and "." in line:
                s = line.split(".", 1)[1].strip()
                if s:
                    steps.append({"title": s[:20], "description": s, "goal": "掌握相关知识", "resource_types": ["文档", "练习"]})

    if not steps:
        weak = ", ".join(profile.get("weak_points", [])) or "无"
        steps = [
            {"title": "夯实基础", "description": f"复习核心概念和基本原理，重点关注{weak}", "goal": "建立知识框架", "resource_types": ["文档", "视频"]},
            {"title": "专项练习", "description": "针对薄弱点进行针对性练习和巩固", "goal": "攻克薄弱环节", "resource_types": ["练习", "案例"]},
            {"title": "综合实践", "description": "完成综合案例，将知识融会贯通", "goal": "综合运用", "resource_types": ["案例", "练习"]},
        ]

    for i, step in enumerate(steps):
        step["status"] = "current" if i == 0 else "todo"

    # 5. 持久化路径
    try:
        db_save_plan(session_id, diagnosis, summary, steps)
    except Exception:
        pass

    return {
        "success": True,
        "plan": {"steps": steps, "summary": summary, "diagnosis": diagnosis},
        "profile": profile,
    }


@app.post("/api/doubao-image")
async def doubao_image(request: Request):
    """豆包(Seedream)图片生成代理端点。"""
    try:
        data = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="无效的 JSON 请求体")
    prompt = data.get("prompt", "").strip()
    if not prompt:
        raise HTTPException(status_code=400, detail="缺少 prompt 参数")
    if len(prompt) > 2000:
        raise HTTPException(status_code=400, detail="prompt 过长（最多2000字符）")

    result = call_doubao_image(prompt)
    return result


@app.post("/api/proxy-download")
async def proxy_download(request: Request):
    """代理下载远程图片，解决跨域问题。支持 JSON 和表单两种格式。"""
    content_type = request.headers.get("content-type", "")
    if "application/json" in content_type:
        data = await request.json()
    else:
        form = await request.form()
        data = {k: v for k, v in form.items()}
    url = data.get("url", "").strip()
    filename = data.get("filename", "download.jpg")
    if not url:
        raise HTTPException(status_code=400, detail="缺少 url 参数")

    try:
        resp = requests.get(url, timeout=30)
        resp.raise_for_status()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"下载图片失败: {e}")

    safe_fn = _re.sub(r'[\\/*?:"<>|]', "_", filename)[:80]
    encoded_fn = quote(safe_fn)
    return Response(
        content=resp.content,
        media_type=resp.headers.get("Content-Type", "image/jpeg"),
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{encoded_fn}"}
    )


@app.post("/api/render-image")
async def render_mermaid_image(request: Request):
    """服务端渲染 Mermaid 代码为 JPG 图片。"""
    data = await request.json()
    code = data.get("code", "").strip()
    if not code:
        raise HTTPException(status_code=400, detail="缺少 mermaid 代码")
    if len(code) > 10000:
        raise HTTPException(status_code=400, detail="代码过长")

    script = os.path.join(os.path.dirname(__file__), "render_mermaid.cjs")
    try:
        proc = await asyncio.create_subprocess_exec(
            "node", script,
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(
            proc.communicate(input=code.encode("utf-8")), timeout=30
        )
        if proc.returncode != 0:
            err_msg = stderr.decode("utf-8", errors="replace")[:200]
            raise HTTPException(status_code=500, detail=f"渲染失败: {err_msg}")
        svg = stdout.decode("utf-8", errors="replace")

        # SVG → PNG → JPG 转换
        import cairosvg
        from PIL import Image

        png_data = cairosvg.svg2png(bytestring=svg.encode("utf-8"))
        img = Image.open(_io.BytesIO(png_data))
        if img.mode in ("RGBA", "P"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode == "RGBA":
                bg.paste(img, mask=img.split()[3])
            else:
                bg.paste(img)
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")

        jpg_buf = _io.BytesIO()
        img.save(jpg_buf, format="JPEG", quality=90)
        jpg_buf.seek(0)
        return Response(content=jpg_buf.getvalue(), media_type="image/jpeg")
    except asyncio.TimeoutError:
        raise HTTPException(status_code=504, detail="渲染超时")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=8000, reload=True, reload_dirs=[os.path.dirname(__file__)])
